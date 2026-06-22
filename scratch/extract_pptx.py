import pptx
import os

def extract_text_from_pptx(pptx_path):
    if not os.path.exists(pptx_path):
        print(f"File not found: {pptx_path}")
        return
    
    prs = pptx.Presentation(pptx_path)
    text_runs = []
    for i, slide in enumerate(prs.slides):
        text_runs.append(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    
    return "\n".join(text_runs)

pptx_file = "cours1_and_av_ing3_DS_2026 (1).pptx"
try:
    content = extract_text_from_pptx(pptx_file)
    with open("pptx_content.txt", "w", encoding="utf-8") as f:
        f.write(content)
    print("Content extracted successfully to pptx_content.txt")
except Exception as e:
    print(f"Error: {e}")
